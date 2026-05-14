import asyncio
import os
import json
import re
import requests
import io
from playwright.async_api import async_playwright
from bs4 import BeautifulSoup
import pdfplumber


class NetologyScraper:
    def __init__(self, cookies_file="data/netology_cookies.json"):
        self.cookies_file = cookies_file
        self.browser = None
        self.context = None
        self.page = None
        self.playwright = None

    async def start(self):
        self.playwright = await async_playwright().start()
        self.browser = await self.playwright.chromium.launch(
            headless=False,
            args=[
                "--no-sandbox",
                "--disable-setuid-sandbox",
                "--disable-dev-shm-usage",
                "--disable-gpu",
            ],
        )
        self.context = await self.browser.new_context()
        if os.path.exists(self.cookies_file):
            with open(self.cookies_file, "r", encoding="utf-8") as f:
                cookies = json.load(f)
            await self.context.add_cookies(cookies)
            print("🍪 Cookies загружены")
        self.page = await self.context.new_page()

    async def stop(self):
        if self.browser:
            await self.browser.close()
        if self.playwright:
            await self.playwright.stop()

    async def save_cookies(self):
        if self.context:
            cookies = await self.context.cookies()
            os.makedirs(os.path.dirname(self.cookies_file) or ".", exist_ok=True)
            with open(self.cookies_file, "w", encoding="utf-8") as f:
                json.dump(cookies, f, ensure_ascii=False, indent=2)
            print("🍪 Cookies сохранены")

    def _ensure_url(self, url):
        if not url:
            return ""
        if not url.startswith("http"):
            return "https://netology.ru" + (url if url.startswith("/") else "/" + url)
        return url

    async def _safe_goto(self, url, wait_until="networkidle", timeout=30000):
        try:
            await self.page.goto(url, wait_until=wait_until, timeout=timeout)
            return True
        except Exception:
            try:
                await self.page.goto(url, wait_until="domcontentloaded", timeout=15000)
                return True
            except Exception:
                return False

    @staticmethod
    def _parse_vtt(vtt_text: str) -> str:
        lines = vtt_text.splitlines()
        result = []
        for line in lines:
            line = line.strip()
            if line.upper() == "WEBVTT":
                continue
            if not line:
                continue
            if re.match(r'^\d{2}:\d{2}:\d{2}\.\d{3}\s*-->', line):
                continue
            if re.match(r'^\d+$', line):
                continue
            if line.upper().startswith(("NOTE", "REGION", "STYLE")):
                continue
            result.append(line)
        return " ".join(result)

    async def _extract_vtt_text(self, url: str) -> str:
        vtt_url = None

        def handle_response(response):
            nonlocal vtt_url
            req_url = response.url
            if ".vtt" in req_url and not vtt_url:
                vtt_url = req_url
                print(f"🎯 VTT: {req_url[:80]}...")

        self.page.on("response", handle_response)
        ok = await self._safe_goto(url, wait_until="domcontentloaded", timeout=60000)
        if not ok:
            self.page.remove_listener("response", handle_response)
            return ""

        for i in range(15):
            if vtt_url:
                break
            await asyncio.sleep(1)
        self.page.remove_listener("response", handle_response)

        if not vtt_url:
            return ""
        try:
            r = requests.get(vtt_url, timeout=30)
            if r.status_code == 200:
                text = self._parse_vtt(r.text)
                print(f"✅ VTT: {len(text)} символов")
                return text
        except Exception as e:
            print(f"⚠️ VTT ошибка: {e}")
        return ""

    async def _extract_file_text(self, url: str) -> str:
        file_url = None
        file_ext = None

        def handle_response(response):
            nonlocal file_url, file_ext
            ct = response.headers.get("content-type", "")
            ul = response.url.lower()
            if not file_url:
                if "pdf" in ct or ul.endswith(".pdf"):
                    file_url, file_ext = response.url, "pdf"
                elif "word" in ct or ul.endswith(".docx"):
                    file_url, file_ext = response.url, "docx"
                elif "powerpoint" in ct or ul.endswith(".pptx"):
                    file_url, file_ext = response.url, "pptx"

        self.page.on("response", handle_response)
        ok = await self._safe_goto(url, wait_until="domcontentloaded", timeout=60000)
        if not ok:
            self.page.remove_listener("response", handle_response)
            return ""

        for i in range(10):
            if file_url:
                break
            await asyncio.sleep(1)
        self.page.remove_listener("response", handle_response)

        if not file_url:
            return ""
        try:
            r = requests.get(file_url, timeout=30)
            if r.status_code != 200:
                return ""
            if file_ext == "pdf":
                with pdfplumber.open(io.BytesIO(r.content)) as pdf:
                    return "\n".join(page.extract_text() or "" for page in pdf.pages)
            elif file_ext == "docx":
                import docx
                doc = docx.Document(io.BytesIO(r.content))
                return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            elif file_ext == "pptx":
                from pptx import Presentation
                prs = Presentation(io.BytesIO(r.content))
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            texts.append(shape.text)
                return "\n".join(texts)
        except Exception as e:
            print(f"⚠️ Файл ошибка: {e}")
        return ""

    async def get_program_disciplines(self, program_id):
        url = f"https://netology.ru/profile/program/{program_id}/schedule"
        print(f"🌐 {url}")
        await self._safe_goto(url)
        await asyncio.sleep(3)

        disciplines = await self.page.evaluate("""
        () => {
            const results = [];
            const seenIds = new Set();
            document.querySelectorAll('[data-lesson-id]').forEach(block => {
                const lessonId = block.getAttribute('data-lesson-id');
                if (!lessonId || seenIds.has(lessonId)) return;
                seenIds.add(lessonId);

                const titleEl = block.querySelector('[data-testid="program-lesson-title"]');
                const title = titleEl ? titleEl.textContent.trim() : '';

                const statusEl = block.querySelector('[data-testid="resourcepack-lesson-status"]');
                const statusText = statusEl ? statusEl.textContent.trim() : '';
                const locked = statusText.toLowerCase().includes('откроется');

                const links = [];
                block.querySelectorAll('a[data-testid="program-granule-link"]').forEach(a => {
                    links.push({text: a.textContent.trim(), href: a.getAttribute('href')});
                });
                if (!links.length) {
                    block.querySelectorAll('a[href*="/lessons/"], a[href*="/lesson_items/"]').forEach(a => {
                        links.push({text: a.textContent.trim(), href: a.getAttribute('href')});
                    });
                }

                results.push({title, lesson_id: lessonId, locked, links});
            });
            return results;
        }
        """)

        raw_title = await self.page.evaluate("""
        () => {
            const el = document.querySelector('[data-testid="program-header"]');
            return el ? el.textContent.trim() : document.title;
        }
        """)

        program_title = re.sub(r'\d+\s+курс.*?:\s*', '', raw_title, flags=re.IGNORECASE)
        program_title = re.sub(r'\d+\s+[а-яА-Я]+\s*—\s*\d+\s+[а-яА-Я]+', '', program_title)
        program_title = re.sub(r'BHEBFAD[-\w]+', '', program_title, flags=re.IGNORECASE)
        program_title = program_title.strip()

        if not disciplines:
            html = await self.page.content()
            debug_path = "data/debug_program.html"
            os.makedirs("data", exist_ok=True)
            with open(debug_path, "w", encoding="utf-8") as f:
                f.write(html)
            print(f"⚠️ Дисциплины не найдены, HTML сохранён: {debug_path}")
            return program_title, []

        print(f"📚 Найдено разделов: {len(disciplines)}")
        for d in disciplines:
            status = "🔒" if d["locked"] else "✅"
            print(f"  {status} {d['title']} (id: {d['lesson_id']})")

        return program_title, disciplines

    async def get_discipline_lessons(self, program_id, lesson_id, fallback_links=None):
        url = f"https://netology.ru/profile/program/{program_id}/lessons/{lesson_id}"
        print(f"🌐 {url}")

        ok = await self._safe_goto(url)
        if not ok:
            print("⚠️ Страница раздела не открылась, пробуем fallback...")
        else:
            await asyncio.sleep(2)
            try:
                await self.page.wait_for_selector('a[data-testid^="program-menu-lessonitem"]', timeout=8000)
            except:
                pass
            await asyncio.sleep(1)

        items = await self.page.evaluate(f"""
        (lessonId) => {{
            const results = [];
            const seen = new Set();
            document.querySelectorAll('a[data-testid^="program-menu-lessonitem"]').forEach(a => {{
                const href = a.getAttribute('href');
                if (!href || seen.has(href)) return;
                if (!href.includes('/lessons/' + lessonId + '/')) return;
                seen.add(href);
                const titleEl = a.querySelector('[data-testid="program-menu-lessonitem-title"]');
                const title = titleEl ? titleEl.textContent.trim() : a.textContent.trim().split('\\n')[0].trim();
                const testid = a.getAttribute('data-testid') || '';
                const locked = testid.toLowerCase().includes('locked');
                results.push({{title, href, locked}});
            }});
            return results;
        }}
        """, lesson_id)

        if not items and fallback_links:
            for link in fallback_links:
                href = link.get("href", "")
                if not href:
                    continue
                print(f"🔁 Fallback: {link.get('text', 'item')}")
                ok = await self._safe_goto(self._ensure_url(href))
                if not ok:
                    continue
                await asyncio.sleep(3)
                try:
                    await self.page.wait_for_selector('a[data-testid^="program-menu-lessonitem"]', timeout=8000)
                except:
                    pass
                await asyncio.sleep(1)

                items = await self.page.evaluate(f"""
                (lessonId) => {{
                    const results = [];
                    const seen = new Set();
                    document.querySelectorAll('a[data-testid^="program-menu-lessonitem"]').forEach(a => {{
                        const href = a.getAttribute('href');
                        if (!href || seen.has(href)) return;
                        if (!href.includes('/lessons/' + lessonId + '/')) return;
                        seen.add(href);
                        const titleEl = a.querySelector('[data-testid="program-menu-lessonitem-title"]');
                        const title = titleEl ? titleEl.textContent.trim() : a.textContent.trim().split('\\n')[0].trim();
                        const testid = a.getAttribute('data-testid') || '';
                        const locked = testid.toLowerCase().includes('locked');
                        results.push({{title, href, locked}});
                    }});
                    return results;
                }}
                """, lesson_id)
                if items:
                    break

        if not items:
            print("⚠️ Занятия не найдены")
            return []

        print(f"📄 Найдено материалов: {len(items)}")
        for item in items:
            status = "🔒" if item["locked"] else "✅"
            print(f"  {status} {item['title']}")
        return items

    async def get_lesson_text_content(self, url):
        url = self._ensure_url(url)
        if not url:
            return ""

        print(f"🌐 {url}")

        vtt_url = None
        file_url = None
        file_ext = None

        def handle_response(response):
            nonlocal vtt_url, file_url, file_ext
            req_url = response.url
            ct = response.headers.get("content-type", "")
            ul = req_url.lower()
            if not vtt_url and ".vtt" in req_url:
                vtt_url = req_url
            if not file_url:
                if "pdf" in ct or ul.endswith(".pdf"):
                    file_url, file_ext = req_url, "pdf"
                elif "word" in ct or ul.endswith(".docx"):
                    file_url, file_ext = req_url, "docx"
                elif "powerpoint" in ct or ul.endswith(".pptx"):
                    file_url, file_ext = req_url, "pptx"

        self.page.on("response", handle_response)

        ok = await self._safe_goto(url, wait_until="domcontentloaded", timeout=60000)
        if not ok:
            self.page.remove_listener("response", handle_response)
            print("⚠️ Страница не открылась")
            return ""

        print("⏳ Ждём загрузку ресурсов...")
        await asyncio.sleep(3)
        self.page.remove_listener("response", handle_response)

        # 0. VTT
        if vtt_url:
            try:
                r = requests.get(vtt_url, timeout=30)
                if r.status_code == 200:
                    text = self._parse_vtt(r.text)
                    if len(text) > 300:
                        print(f"✅ VTT: {len(text)} символов")
                        return text
            except Exception as e:
                print(f"⚠️ VTT ошибка: {e}")

        # 1. Файлы (PDF/DOCX/PPTX)
        if file_url:
            try:
                r = requests.get(file_url, timeout=30)
                if r.status_code != 200:
                    return ""
                if file_ext == "pdf":
                    with pdfplumber.open(io.BytesIO(r.content)) as pdf:
                        text = "\n".join(page.extract_text() or "" for page in pdf.pages)
                elif file_ext == "docx":
                    import docx
                    doc = docx.Document(io.BytesIO(r.content))
                    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                elif file_ext == "pptx":
                    from pptx import Presentation
                    prs = Presentation(io.BytesIO(r.content))
                    texts = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                texts.append(shape.text)
                    text = "\n".join(texts)
                print(f"✅ {file_ext.upper()}: {len(text)} символов")
                if len(text) > 100000:
                    print(f"   ⚠️ Файл слишком большой ({len(text)} символов). Создаём заглушку.")
                    return f"[ФАЙЛ-УЧЕБНИК]\n\nЭтот файл содержит {len(text)} символов и, вероятно, является учебником или дополнительными материалами.\n\nРекомендуется изучить самостоятельно.\n\n---TERMS---\nУчебник|Дополнительный материал для самостоятельного изучения.\n---END_TERMS---"
                if len(text) > 300:
                    return text
            except Exception as e:
                print(f"⚠️ Файл ошибка: {e}")

        # 2. HTML fallback
        html = await self.page.content()
        soup = BeautifulSoup(html, "lxml")

        parts = []
        heading = soup.find("div", attrs={"data-testid": "program-lessonitem-title"})
        if heading:
            parts.append(heading.get_text(strip=True))

        webinar = soup.find("div", attrs={"data-testid": "webinar-info"})
        if webinar:
            parts.append(webinar.get_text(separator="\n", strip=True))

        selectors = [
            'div[data-testid="program-lessonitem-content"]',
            'div[data-testid="lesson-content"]',
            "article",
            "main",
            ".lesson-content",
            "#lesson-content",
            ".content-wrapper",
        ]
        for sel in selectors:
            tag = soup.select_one(sel)
            if tag:
                for junk in tag.find_all(["script", "style", "nav"]):
                    junk.decompose()
                txt = tag.get_text(separator="\n", strip=True)
                if len(txt) > 100:
                    parts.append(txt)
                    break

        result = "\n\n".join(parts)
        if len(result) < 200 and soup.body:
            for junk in soup.body.find_all(["script", "style", "nav", "header", "footer"]):
                junk.decompose()
            body_text = soup.body.get_text(separator="\n", strip=True)
            if len(body_text) > len(result):
                result = body_text[:15000]

        return result[:15000]

    async def debug_webinar(self, url, output_path="data/debug_webinar.html"):
        url = self._ensure_url(url)
        print(f"🔍 Диагностика вебинара: {url}")
        ok = await self._safe_goto(url)
        if not ok:
            print("❌ Не удалось открыть")
            return
        await asyncio.sleep(3)

        html = await self.page.content()
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(html)
        print(f"💾 HTML сохранён: {output_path}")
