"""
Netology API Client — HTTP-first доступ к контенту Нетологии.

Использует httpx + cookies (полученные через Playwright).
Playwright нужен ТОЛЬКО для авторизации.

Endpoints:
  GET /backend/api/user/programs/{id}/schedule  → lessons[] + lesson_items[]
  GET /backend/api/user/lesson_items/{id}       → content / video_url / files[]
"""

import asyncio
import json
import os
import re
import time
from typing import Optional

import httpx
import requests

API_BASE = "https://netology.ru/backend/api/user"


class NetologyAPIClient:
    """HTTP-клиент для API Нетологии."""

    def __init__(self, cookies_file: str = "backend/netology_cookies/netology_cookies.json"):
        self.cookies_file = cookies_file
        self.cookies = self._load_cookies()
        self.client = httpx.AsyncClient(
            cookies=self.cookies,
            follow_redirects=True,
            timeout=httpx.Timeout(30.0),
        )

    def _load_cookies(self) -> dict:
        if not os.path.exists(self.cookies_file):
            return {}
        with open(self.cookies_file, "r", encoding="utf-8") as f:
            raw = json.load(f)
        # httpx принимает dict name -> value
        return {c["name"]: c["value"] for c in raw if "name" in c}

    async def close(self):
        await self.client.aclose()

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    async def get_program_info(self, program_id: str | int):
        """
        Возвращает метаданные программы (название, даты и т.д.).
        Полезно для получения program_title.
        """
        url = f"{API_BASE}/programs/{program_id}"
        r = await self.client.get(url)
        r.raise_for_status()
        return r.json()

    async def get_program_schedule(self, program_id: str | int):
        """
        Возвращает расписание программы (список занятий с lesson_items).
        program_id может быть slug (bhebfad-25-memeo-2) или числом.
        """
        url = f"{API_BASE}/programs/{program_id}/schedule"
        r = await self.client.get(url)
        r.raise_for_status()
        data = r.json()
        return data.get("lessons", [])

    async def get_lesson_item(self, item_id: int | str):
        """
        Возвращает детали lesson_item.
        Для type=text: content (markdown)
        Для type=video/webinar: video_url, webinar_url, youtube_video_id
        Для type=attachment: files[]
        """
        url = f"{API_BASE}/lesson_items/{item_id}"
        r = await self.client.get(url)
        r.raise_for_status()
        return r.json()

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    @staticmethod
    def download_file(file_url: str, timeout: int = 60) -> bytes:
        """Синхронная загрузка файла (PDF/DOCX/PPTX)."""
        r = requests.get(file_url, timeout=timeout)
        r.raise_for_status()
        return r.content

    @staticmethod
    def extract_text_from_pdf_bytes(pdf_bytes: bytes) -> str:
        """Извлекает текст из PDF (байты)."""
        import io
        import pdfplumber
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            return "\n".join(page.extract_text() or "" for page in pdf.pages)

    @staticmethod
    def extract_text_from_docx_bytes(docx_bytes: bytes) -> str:
        """Извлекает текст из DOCX (байты)."""
        import io
        import docx
        doc = docx.Document(io.BytesIO(docx_bytes))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    @staticmethod
    def extract_text_from_pptx_bytes(pptx_bytes: bytes) -> str:
        """Извлекает текст из PPTX (байты)."""
        import io
        from pptx import Presentation
        prs = Presentation(io.BytesIO(pptx_bytes))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n".join(texts)

    @staticmethod
    def parse_file_text(file_bytes: bytes, extension: str) -> str:
        """Универсальный парсер файла по расширению."""
        ext = extension.lower().lstrip(".")
        if ext == "pdf":
            return NetologyAPIClient.extract_text_from_pdf_bytes(file_bytes)
        if ext == "docx":
            return NetologyAPIClient.extract_text_from_docx_bytes(file_bytes)
        if ext == "pptx":
            return NetologyAPIClient.extract_text_from_pptx_bytes(file_bytes)
        return ""

    # ------------------------------------------------------------------
    # High-level: собирает контент item'а целиком
    # ------------------------------------------------------------------

    async def fetch_item_content(self, item: dict, subject_name: str = "") -> dict:
        """
        Получает полный контент item'а:
          - text (str) — извлечённый текст (markdown, PDF, VTT или транскрипция)
          - video_url (str) — если есть видео
          - title, href, locked
          - is_structure (bool) — если текст — описание программы

        Возвращает dict совместимый с текущим _get_item_content.
        """
        item_id = item["id"]
        item_title = item.get("title", "")
        item_type = item.get("type", "")
        locked = item.get("locked", False)
        href = f"https://netology.ru{item['path']}" if item.get("path") else ""

        text = ""
        video_url = ""
        is_structure = False

        # Для text/video/webinar/attachment — делаем запрос к item
        # Но для оптимизации: если type=text и нам нужен только preview,
        #   можно пропустить. Но пока всегда запрашиваем.
        try:
            detail = await self.get_lesson_item(item_id)
        except httpx.HTTPStatusError as e:
            if e.response.status_code == 401:
                raise RuntimeError("Cookies протухли, требуется повторная авторизация") from e
            # 404 или другая ошибка — пропускаем item
            print(f"   ⚠️ Ошибка API item {item_id}: {e.response.status_code}")
            return {
                "title": item_title,
                "href": href,
                "text": "",
                "video_url": "",
                "is_structure": False,
                "locked": locked,
            }

        # Собираем video_url из любых источников
        video_url = detail.get("video_url") or ""
        if not video_url:
            video_url = detail.get("webinar_url") or ""
        if not video_url:
            youtube_id = detail.get("youtube_video_id")
            if youtube_id:
                video_url = f"https://www.youtube.com/watch?v={youtube_id}"

        # Обработка по типу
        if item_type == "text":
            text = detail.get("content", "") or ""
            # Structure page detection (тот же guard, что и раньше)
            if text and NetologyAPIClient._is_structure_page(text):
                is_structure = True

        elif item_type == "attachment":
            files = detail.get("files", [])
            if files:
                file_info = files[0]
                file_url = file_info.get("link", "")
                file_ext = file_info.get("extension", "")
                if file_url:
                    try:
                        file_bytes = self.download_file(file_url)
                        text = self.parse_file_text(file_bytes, file_ext)
                        print(f"   ✅ {file_ext.upper()}: {len(text)} символов")
                    except Exception as e:
                        print(f"   ⚠️ Ошибка загрузки файла: {e}")

        elif item_type in ("video", "webinar"):
            # Контента нет, но video_url уже установлен выше
            pass

        return {
            "title": item_title,
            "href": href,
            "text": text,
            "video_url": video_url,
            "is_structure": is_structure,
            "locked": locked,
        }

    # ------------------------------------------------------------------
    # Static helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _is_structure_page(text: str) -> bool:
        """Определяет, является ли текст страницей-описанием курса."""
        # Импортируем здесь, чтобы избежать циклических зависимостей
        try:
            from material_classifier import is_structure_page as _isp
            return _isp(text)
        except Exception:
            # Fallback если material_classifier недоступен
            indicators = [
                "цель изучения",
                "задачи изучения",
                "содержание дисциплины",
                "тематический план",
                "формы контроля",
                "балльно-рейтинговая система",
                "рекомендуемая литература",
            ]
            t = text.lower()
            score = sum(1 for ind in indicators if ind in t)
            return score >= 3 and len(text) < 5000
